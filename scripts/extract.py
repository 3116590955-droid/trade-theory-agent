import sys
sys.stdout.reconfigure(encoding='utf-8')

from pathlib import Path
import fitz  # pymupdf
from pptx import Presentation

SOURCE_DIR = Path(__file__).parent.parent / "国际经济学课上知识"
OUTPUT_DIR = Path(__file__).parent.parent / "knowledge_base"


def extract_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def write_markdown(src: Path, out_dir: Path) -> None:
    suffix = src.suffix.lower()
    if suffix == ".pdf":
        body = extract_pdf(src)
        file_type = "pdf"
    elif suffix == ".pptx":
        body = extract_pptx(src)
        file_type = "pptx"
    else:
        return  # 跳过不支持的格式

    header = f"---\nsource: {src.name}\ntype: {file_type}\noriginal_path: 国际经济学课上知识/{src.name}\n---\n\n"
    out_path = out_dir / (src.stem + ".md")
    out_path.write_text(header + body, encoding="utf-8")
    print(f"  OK {src.name} -> {out_path.name}")


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)
    files = sorted(SOURCE_DIR.iterdir())
    print(f"Processing {len(files)} files...")
    for f in files:
        write_markdown(f, OUTPUT_DIR)
    print("Done.")


if __name__ == "__main__":
    main()
